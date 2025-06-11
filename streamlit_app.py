import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
import csv
import os
import io  # Required for in-memory file handling
import traceback


# --- Helper Functions (Copied from your original script, no changes needed) ---

def get_shape_fill_color_info(shape):
    """
    Determines the fill color of a shape and returns a descriptive string.
    """
    if not hasattr(shape, 'fill'):
        if hasattr(MSO_SHAPE_TYPE, 'TABLE') and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "Table Container (Fill N/A)"
        elif hasattr(MSO_SHAPE_TYPE, 'CHART') and shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return "Chart Container (Fill N/A)"
        is_other_graphic_frame = (hasattr(MSO_SHAPE_TYPE,
                                          'GRAPHIC_FRAME') and shape.shape_type == MSO_SHAPE_TYPE.GRAPHIC_FRAME) or \
                                 (shape.shape_type == 6)
        if is_other_graphic_frame:
            return "Graphic Frame (Fill N/A)"
        return "Fill Attribute Missing / Other Type"
    fill = shape.fill
    color_info = "No Fill"
    try:
        fill_type_map = {
            MSO_FILL_TYPE.SOLID: "Solid Fill", MSO_FILL_TYPE.GRADIENT: "Gradient Fill",
            MSO_FILL_TYPE.PICTURE: "Picture Fill", MSO_FILL_TYPE.GROUP: "Group Fill",
            MSO_FILL_TYPE.BACKGROUND: "Background Fill", MSO_FILL_TYPE.PATTERNED: "Patterned Fill",
            MSO_FILL_TYPE.TEXTURED: "Textured Fill"
        }
        if fill.type in fill_type_map:
            color_info = fill_type_map[fill.type]
            if fill.type == MSO_FILL_TYPE.SOLID:
                color = fill.fore_color
                if hasattr(color, 'rgb') and color.rgb is not None:
                    rgb_val = color.rgb
                    color_info = f"RGB({rgb_val[0]},{rgb_val[1]},{rgb_val[2]})"
                elif hasattr(color, 'theme_color') and color.theme_color is not None:
                    brightness = color.brightness
                    color_info = f"ThemeColor(Type:{color.theme_color}, Brightness:{brightness:.2f})"
    except Exception:
        color_info = "Fill Info Error (General)"
    return color_info


def get_all_text_from_shape(shape):
    """
    Extracts and concatenates all text from a shape, whether it's in a text frame or a table.
    """
    text_parts = []
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text.strip())
    is_table = (hasattr(MSO_SHAPE_TYPE, 'TABLE') and shape.shape_type == MSO_SHAPE_TYPE.TABLE) or \
               (not hasattr(MSO_SHAPE_TYPE, 'TABLE') and shape.shape_type == 19)
    if is_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    text_parts.append(cell.text.strip())
    return " | ".join(text_parts)


# --- Core Processing Functions (Revised for Streamlit) ---

def generate_text_details_tsv(uploaded_file):
    """
    MODIFIED: Loads a PowerPoint file from an in-memory object, extracts detailed text data,
    and returns it as a string in TSV format.
    """
    output = io.StringIO()  # Create an in-memory text buffer
    tsv_writer = csv.writer(output, delimiter='\t')

    try:
        prs = Presentation(uploaded_file)
        headers = ["slide_index", "shape_id", "is_table_cell", "cell_row_index", "cell_col_index",
                   "paragraph_index", "run_index", "text"]
        tsv_writer.writerow(headers)

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if (hasattr(MSO_SHAPE_TYPE, 'LINE') and shape.shape_type == MSO_SHAPE_TYPE.LINE) or \
                        (not hasattr(MSO_SHAPE_TYPE, 'LINE') and shape.shape_type == 13):
                    continue

                shape_id = shape.shape_id
                if shape.has_text_frame:
                    for para_idx, p in enumerate(shape.text_frame.paragraphs):
                        for run_idx, run in enumerate(p.runs):
                            if run.text.strip():
                                tsv_writer.writerow(
                                    [slide_idx, shape_id, False, None, None, para_idx, run_idx, run.text.strip()])
                if (hasattr(MSO_SHAPE_TYPE, 'TABLE') and shape.shape_type == MSO_SHAPE_TYPE.TABLE) or \
                        (not hasattr(MSO_SHAPE_TYPE, 'TABLE') and shape.shape_type == 19):
                    for r_idx, row in enumerate(shape.table.rows):
                        for c_idx, cell in enumerate(row.cells):
                            for para_idx, p in enumerate(cell.text_frame.paragraphs):
                                for run_idx, run in enumerate(p.runs):
                                    if run.text.strip():
                                        tsv_writer.writerow([slide_idx, shape_id, True, r_idx, c_idx, para_idx, run_idx,
                                                             run.text.strip()])
    except Exception as e:
        st.error(f"An error occurred while extracting text details: {e}")
        traceback.print_exc()
        return None

    return output.getvalue()


def generate_combined_shape_details_tsv(uploaded_file):
    """
    MODIFIED: Loads a PowerPoint file from an in-memory object, extracts shape details,
    and returns it as a string in TSV format.
    """
    output = io.StringIO()  # Create an in-memory text buffer
    csv_headers = ["slide_index", "shape_id", "shape_name", "shape_type", "color",
                   "x_coordinate_emu", "y_coordinate_emu", "width_emu", "height_emu", "text"]
    tsv_writer = csv.DictWriter(output, fieldnames=csv_headers, delimiter='\t')
    tsv_writer.writeheader()

    shape_details_list = []
    try:
        # IMPORTANT: Reset read position of the file-like object for the second parse
        uploaded_file.seek(0)
        prs = Presentation(uploaded_file)

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if (hasattr(MSO_SHAPE_TYPE, 'LINE') and shape.shape_type == MSO_SHAPE_TYPE.LINE) or \
                        (not hasattr(MSO_SHAPE_TYPE, 'LINE') and shape.shape_type == 13):
                    continue

                shape_type_str = next((m.name for m in MSO_SHAPE_TYPE if m.value == shape.shape_type),
                                      f"Unknown ({shape.shape_type})")

                shape_details_list.append({
                    "slide_index": slide_idx, "shape_id": shape.shape_id,
                    "shape_name": shape.name or "Unnamed Shape", "shape_type": shape_type_str,
                    "color": get_shape_fill_color_info(shape),
                    "x_coordinate_emu": shape.left if hasattr(shape, 'left') else 'N/A',
                    "y_coordinate_emu": shape.top if hasattr(shape, 'top') else 'N/A',
                    "width_emu": shape.width if hasattr(shape, 'width') else 'N/A',
                    "height_emu": shape.height if hasattr(shape, 'height') else 'N/A',
                    "text": get_all_text_from_shape(shape)
                })

        tsv_writer.writerows(shape_details_list)

    except Exception as e:
        st.error(f"An error occurred while extracting shape details: {e}")
        traceback.print_exc()
        return None

    return output.getvalue()


# --- Streamlit User Interface ---

st.set_page_config(page_title="PPTX Element Extractor", layout="centered")
st.title("PowerPoint Element Extractor")
st.write(
    "Upload a PowerPoint (.pptx) file to extract detailed information about its text and shapes into downloadable TSV files. TSV stands for tab separated values. These files are similar to csv files but use a tab as the separator instead of a comma.")

uploaded_file = st.file_uploader("Choose a .pptx file", type="pptx")

if uploaded_file is not None:
    st.success(f"File '{uploaded_file.name}' uploaded successfully!")

    with st.spinner('Processing your presentation... This may take a moment.'):
        # Generate both TSV files in memory
        text_tsv_content = generate_text_details_tsv(uploaded_file)
        shape_tsv_content = generate_combined_shape_details_tsv(uploaded_file)

    st.info("Processing complete. Your files are ready for download below.")

    base_filename = os.path.splitext(uploaded_file.name)[0]

    # Display download buttons if content was generated successfully
    if text_tsv_content:
        st.download_button(
            label="Download Text Details (TSV)",
            data=text_tsv_content,
            file_name=f"{base_filename}_TextDetails.tsv",
            mime="text/tab-separated-values"
        )

    if shape_tsv_content:
        st.download_button(
            label="Download Shape Summary (TSV)",
            data=shape_tsv_content,
            file_name=f"{base_filename}_ShapeSummary.tsv",
            mime="text/tab-separated-values"
        )